"""Studio 生成器（NotebookLM 式)——从知识库源用 LLM 生成多种产物。

文本类(qwen3.7-plus):study_guide / faq / briefing / timeline(Markdown);
结构类(JSON):mind_map(节点/边)/ flashcards(问答)/ quiz(选择题)/ data_table(表)。
音频/幻灯/视频类在 studio_media(TTS/渲染)。
"""

from __future__ import annotations

import json
import re
import uuid

import structlog
from sqlalchemy.ext.asyncio import AsyncSession

from app.services import model_client
from app.services.model_client import ModelError

log = structlog.get_logger("terrane.studio")

_MAX = 12000

# kind → (是否JSON, 提示词)
_PROMPTS: dict[str, tuple[bool, str]] = {
    "study_guide": (False,
        "把下面资料整理成一份结构化「学习指南」(Markdown):## 概述、## 关键概念(要点)、## 重点难点、## 复习问题。只基于资料,简洁准确。"),
    "faq": (False,
        "基于下面资料生成一份「常见问答 FAQ」(Markdown),8-12 组,每组 **问:** 与 答。只基于资料,不编造。"),
    "briefing": (False,
        "基于下面资料写一份「简报文档」(Markdown):一段执行摘要 + 3-6 个要点小节。客观、精炼。"),
    "timeline": (False,
        "从下面资料抽取「时间线」(Markdown 有序列表,每项「时间 — 事件」)。无明确时间则按逻辑顺序。只基于资料。"),
    "mind_map": (True,
        "从下面资料生成思维导图,只输出 JSON:{\"root\":\"中心主题\",\"nodes\":[{\"id\":\"n1\",\"label\":\"分支\",\"parent\":\"root\"},...]}。层级清晰,parent 用上级 id 或 root。"),
    "flashcards": (True,
        "从下面资料生成记忆闪卡,只输出 JSON 数组 [{\"front\":\"问题/概念\",\"back\":\"答案/解释\"}],10-15 张,只基于资料。"),
    "quiz": (True,
        "从下面资料生成测验,只输出 JSON 数组 [{\"q\":\"题干\",\"options\":[\"A\",\"B\",\"C\",\"D\"],\"answer\":0,\"explain\":\"解析\"}],answer 为正确项下标。8-10 题,只基于资料。"),
    "data_table": (True,
        "从下面资料抽取结构化数据,只输出 JSON {\"columns\":[\"列1\",\"列2\",...],\"rows\":[[...],...]}。挑最有信息量的维度,只基于资料。"),
    "slide_deck": (True,
        "从下面资料生成演示文稿大纲,只输出 JSON {\"title\":\"演示标题\",\"subtitle\":\"副标题\",\"slides\":[{\"title\":\"页标题\",\"bullets\":[\"要点1\",\"要点2\",...]},...]}。8-14 页,每页 3-5 个要点,逻辑清晰,只基于资料。"),
}

KINDS = tuple(_PROMPTS)


def _extract_json(raw: str):
    raw = raw.strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```[a-zA-Z]*\n?|\n?```$", "", raw).strip()
    m = re.search(r"[\[{].*[\]}]", raw, re.S)
    return json.loads(m.group(0)) if m else None


async def generate(db: AsyncSession, *, kind: str, sources: list[tuple[str, str]]) -> dict:
    """生成某类 Studio 产物。返回 {kind, format, content}。"""
    if kind not in _PROMPTS:
        raise ValueError(f"unknown studio kind: {kind}")
    is_json, instruction = _PROMPTS[kind]
    corpus = "\n\n".join(f"## {t}\n{x}" for t, x in sources if x and x.strip())[:_MAX]
    if not corpus.strip():
        return {"kind": kind, "format": "empty", "content": None}
    sys = "你是知识库 Studio 生成器。严格基于给定资料,不编造。" + ("只输出 JSON,无任何多余文字。" if is_json else "")
    raw = await model_client.chat_complete(
        db, [{"role": "system", "content": sys},
             {"role": "user", "content": instruction + "\n\n【资料】\n" + corpus}],
        temperature=0.3, max_tokens=3000)
    if is_json:
        try:
            data = _extract_json(raw)
        except (ValueError, AttributeError):
            data = None
        if data is None:
            raise ModelError("studio_json_parse_failed")
        return {"kind": kind, "format": "json", "content": data}
    return {"kind": kind, "format": "markdown", "content": raw.strip()}


# ---- 媒体类:幻灯片(PPTX)/ 播客音频(TTS)----

def build_pptx(deck: dict) -> bytes:
    """从 slide_deck JSON 生成 .pptx 字节。"""
    import io

    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = str(deck.get("title") or "演示文稿")
    if deck.get("subtitle") and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = str(deck["subtitle"])
    for sl in (deck.get("slides") or [])[:30]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(sl.get("title") or "")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = sl.get("bullets") or []
        for i, b in enumerate(bullets[:8]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(b)
            p.font.size = Pt(18)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _concat_wav(segs: list[bytes]) -> bytes:
    """拼接多段 wav 为一段(同采样参数)。非 PCM wav 则原样首段返回。"""
    import io
    import wave

    valid = [s for s in segs if s]
    if not valid:
        return b""
    try:
        out = io.BytesIO()
        writer = None
        for s in valid:
            rf = wave.open(io.BytesIO(s), "rb")
            if writer is None:
                writer = wave.open(out, "wb")
                writer.setnchannels(rf.getnchannels())
                writer.setsampwidth(rf.getsampwidth())
                writer.setframerate(rf.getframerate())
            writer.writeframes(rf.readframes(rf.getnframes()))
            rf.close()
        if writer:
            writer.close()
        return out.getvalue()
    except (wave.Error, EOFError):
        return valid[0]


async def generate_podcast(db: AsyncSession, *, sources: list[tuple[str, str]]) -> dict:
    """双人播客:生成对话脚本 + TTS 合成,返回 {script, audio(data url)}。无 TTS 渠道→error。"""
    corpus = "\n\n".join(f"## {t}\n{x}" for t, x in sources if x and x.strip())[:_MAX]
    if not corpus.strip():
        return {"kind": "audio_overview", "format": "empty", "content": None}
    raw = await model_client.chat_complete(
        db, [{"role": "system", "content": "你是播客编剧。只输出 JSON,无多余文字。"},
             {"role": "user", "content":
              "把资料改写成双人中文播客对话,只输出 JSON 数组 [{\"speaker\":\"A\"|\"B\",\"text\":\"...\"}]。"
              "A=主持人(提问/引导/串场),B=专家(讲解)。12-18 轮,口语化、自然、信息准确,只基于资料。\n\n【资料】\n" + corpus}],
        temperature=0.5, max_tokens=3000)
    try:
        lines = _extract_json(raw)
    except (ValueError, AttributeError):
        lines = None
    if not isinstance(lines, list) or not lines:
        raise ModelError("podcast_script_failed")
    import asyncio

    voices = {"A": "Cherry", "B": "Ethan"}
    segs: list[bytes] = []
    used: list[dict] = []
    for i, ln in enumerate(lines[:12]):
        text = str(ln.get("text") or "").strip()
        if not text:
            continue
        if i:
            await asyncio.sleep(0.8)  # 节流,避开 qwen-tts QPS 限流
        audio = await model_client.tts(db, text, voice=voices.get(str(ln.get("speaker")), "Cherry"))
        if audio is None:
            raise ModelError("no_tts_channel")
        segs.append(audio)
        used.append({"speaker": ln.get("speaker"), "text": text})
    import base64
    wav = _concat_wav(segs)
    return {"kind": "audio_overview", "format": "audio",
            "content": {"script": used, "audio": "data:audio/wav;base64," + base64.b64encode(wav).decode()}}
