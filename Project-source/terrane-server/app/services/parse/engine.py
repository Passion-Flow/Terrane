"""Terrane Parse -- in-house document parsing engine (pure CPU, no GPU). See design/07-parse-engine.md.

Mature libraries (pdfplumber / python-docx ... = the text-extraction layer) provide the low-level primitives; the "intelligence" -- layout reconstruction, table reconstruction, etc. -- is 100% in-house:
- Reading order: column detection (vertical whitespace) + top->down within a column (XY-Cut idea).
- Headings: graded by font size relative to body text.
- Tables: ruled -> intersect a grid from vector lines; unruled -> cluster text coordinates (this version does ruled tables first; unruled is an enhancement).
Outputs structured Markdown and feeds into the existing chunking / embedding / graph pipeline.
"""

from __future__ import annotations

import re
import statistics

import structlog

log = structlog.get_logger("terrane.parse")

# Formula detection: math font OR a high density of Unicode math symbols. In-house, conservative (avoids misclassifying body text).
_MATH_FONT = re.compile(r"math|cmmi|cmsy|cmex|stix|mathjax|msam|msbm", re.I)
# Only accept unambiguous math symbols (excluding common or ambiguous characters like ·/°/×, to avoid misclassifying body text / placeholders).
_MATH_CHARS = set("∑∫√∞≤≥≠≈≅∀∃∈∉⊂⊆⊃∪∩∂∇∆∏≡∝⊥∠⇒⇔↦∮∭∬")


def _line_is_formula(spans: list) -> bool:
    if any(_MATH_FONT.search(s.get("font", "")) for s in spans):
        return True
    stripped = "".join(s.get("text", "") for s in spans).strip()
    if len(stripped) < 2:
        return False
    m = sum(1 for c in stripped if c in _MATH_CHARS)
    return m >= 3 and m / len(stripped) > 0.2   # Strict threshold: at least 3 strong math symbols and a high density

SUPPORTED = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}


def _cluster(vals: list[float], tol: float = 3.0) -> list[float]:
    """Cluster nearby coordinates into a single line (returns each cluster's mean, ascending)."""
    if not vals:
        return []
    vs = sorted(vals)
    clusters = [[vs[0]]]
    for v in vs[1:]:
        if v - clusters[-1][-1] <= tol:
            clusters[-1].append(v)
        else:
            clusters.append([v])
    return [statistics.mean(c) for c in clusters]


def _ruled_tables(page, words_xy: list[tuple]) -> list[dict]:
    """Reconstruct ruled tables from vector lines. Returns [{bbox, md}]. In-house: lines -> grid -> cells -> placement.

    `page` is a pdfplumber Page; `words_xy` is the page's words as (x0,y0,x1,y1,text) tuples (top-origin y).
    """
    hlines: list[tuple[float, float, float]] = []  # (y, x0, x1)
    vlines: list[tuple[float, float, float]] = []  # (x, y0, y1)
    for ln in page.lines:
        ax0, ay0, ax1, ay1 = ln["x0"], ln["top"], ln["x1"], ln["bottom"]
        if abs(ay0 - ay1) <= 2 and abs(ax0 - ax1) > 8:
            hlines.append((ay0, min(ax0, ax1), max(ax0, ax1)))
        elif abs(ax0 - ax1) <= 2 and abs(ay0 - ay1) > 8:
            vlines.append((ax0, min(ay0, ay1), max(ay0, ay1)))
    for r in page.rects:
        rx0, ry0, rx1, ry1 = r["x0"], r["top"], r["x1"], r["bottom"]
        hlines += [(ry0, rx0, rx1), (ry1, rx0, rx1)]
        vlines += [(rx0, ry0, ry1), (rx1, ry0, ry1)]
    if len(hlines) < 2 or len(vlines) < 2:
        return []
    ys = _cluster([h[0] for h in hlines])
    xs = _cluster([v[0] for v in vlines])
    if len(ys) < 2 or len(xs) < 2:
        return []
    x0, x1, y0, y1 = xs[0], xs[-1], ys[0], ys[-1]
    # Take text words within the table region and place them into cells
    words = words_xy  # (x0,y0,x1,y1,word)
    rows_md: list[list[str]] = []
    for ri in range(len(ys) - 1):
        row: list[str] = []
        for ci in range(len(xs) - 1):
            cx0, cx1, cy0, cy1 = xs[ci], xs[ci + 1], ys[ri], ys[ri + 1]
            cell_words = [w[4] for w in words
                          if cx0 - 1 <= (w[0] + w[2]) / 2 <= cx1 + 1 and cy0 - 1 <= (w[1] + w[3]) / 2 <= cy1 + 1]
            row.append(" ".join(cell_words).strip())
        rows_md.append(row)
    if not rows_md or not any(any(c for c in r) for r in rows_md):
        return []
    md = _rows_to_md(rows_md)
    return [{"bbox": (x0, y0, x1, y1), "md": md}]


def _rows_to_md(rows: list[list[str]]) -> str:
    ncol = max(len(r) for r in rows)
    rows = [r + [""] * (ncol - len(r)) for r in rows]
    out = ["| " + " | ".join(c.replace("|", "\\|") or " " for c in rows[0]) + " |",
           "| " + " | ".join(["---"] * ncol) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(c.replace("|", "\\|") or " " for c in r) + " |")
    return "\n".join(out)


_COL_GAP = 18.0   # Minimum horizontal gap between columns (pt)


def _borderless_regions(blocks: list[dict]) -> list[dict]:
    """Page-level unruled table reconstruction (in-house): cluster all page spans into rows by y, split each row into cells by x gaps,
    and a run of >=2 consecutive multi-column rows with a consistent column count -> one table. Conservative, guards strictly against misclassifying body text. Returns [{bbox, md}]."""
    spans = [s for b in blocks for ln in b.get("lines", []) for s in ln.get("spans", [])
             if s.get("text", "").strip()]
    if len(spans) < 4:
        return []
    spans.sort(key=lambda s: (s["bbox"][1] + s["bbox"][3]) / 2)
    # Cluster into rows
    rows: list[list] = []
    for s in spans:
        yc = (s["bbox"][1] + s["bbox"][3]) / 2
        if rows and abs(yc - (rows[-1][0]["bbox"][1] + rows[-1][0]["bbox"][3]) / 2) < 6:
            rows[-1].append(s)
        else:
            rows.append([s])
    # Row -> cells (by x gaps)
    rc: list[list[tuple[float, str]]] = []
    for r in rows:
        r.sort(key=lambda s: s["bbox"][0])
        cells: list[list] = [[r[0]]]
        for s in r[1:]:
            if s["bbox"][0] - cells[-1][-1]["bbox"][2] > _COL_GAP:
                cells.append([s])
            else:
                cells[-1].append(s)
        rc.append([(c[0]["bbox"][0], " ".join(sp["text"] for sp in c).strip()) for c in cells])

    out: list[dict] = []
    i = 0
    while i < len(rc):
        if len(rc[i]) < 2:
            i += 1
            continue
        # Collect downward the consecutive multi-column rows with the same column count (+/-0)
        j = i
        while j < len(rc) and len(rc[j]) == len(rc[i]) and len(rc[j]) >= 2:
            j += 1
        run = rc[i:j]
        if len(run) >= 2:   # >=2 rows form a table
            xs = _cluster(sorted(x for r in run for x, _ in r), tol=12)
            if len(xs) >= 2:
                grid = []
                for r in run:
                    cols = [""] * len(xs)
                    for x, txt in r:
                        ci = min(range(len(xs)), key=lambda k: abs(xs[k] - x))
                        cols[ci] = (cols[ci] + " " + txt).strip()
                    grid.append(cols)
                ys = [s["bbox"] for rr in rows[i:j] for s in rr]
                bbox = (min(b[0] for b in ys), min(b[1] for b in ys),
                        max(b[2] for b in ys), max(b[3] for b in ys))
                out.append({"bbox": bbox, "md": _rows_to_md(grid)})
        i = max(j, i + 1)
    return out


def _inside(bbox, region, pad: float = 2.0) -> bool:
    bx0, by0, bx1, by1 = bbox
    rx0, ry0, rx1, ry1 = region
    cx, cy = (bx0 + bx1) / 2, (by0 + by1) / 2
    return rx0 - pad <= cx <= rx1 + pad and ry0 - pad <= cy <= ry1 + pad


def _reading_order(blocks: list[dict], page_w: float) -> list[dict]:
    """Column detection (simplified XY-Cut): split into 1-2 columns by x center, top->down within a column."""
    if not blocks:
        return []
    mid = page_w / 2
    left = [b for b in blocks if (b["bbox"][0] + b["bbox"][2]) / 2 < mid]
    right = [b for b in blocks if (b["bbox"][0] + b["bbox"][2]) / 2 >= mid]
    # Treat as two columns only when both columns have substantial content and few blocks straddle the divide
    straddle = [b for b in blocks if b["bbox"][0] < mid - page_w * 0.08 and b["bbox"][2] > mid + page_w * 0.08]
    if len(left) >= 2 and len(right) >= 2 and len(straddle) <= len(blocks) * 0.2:
        return sorted(left, key=lambda b: b["bbox"][1]) + sorted(right, key=lambda b: b["bbox"][1])
    return sorted(blocks, key=lambda b: b["bbox"][1])


_BOLD_FLAG = 1 << 4   # bold flag bit the downstream heading heuristics check (s["flags"] & 16)


def _page_to_blocks(page) -> list[dict]:
    """Convert a pdfplumber page into structured text blocks (one line per text line, spans = words).

    Each span carries bbox (top-origin y), text, size and font so the in-house heading / table / formula
    heuristics keep working unchanged."""
    try:
        words = page.extract_words(extra_attrs=["size", "fontname"])
    except Exception:  # noqa: BLE001
        return []
    words = [w for w in words if (w.get("text") or "").strip()]
    if not words:
        return []
    words.sort(key=lambda w: (round(float(w.get("top", 0.0)), 1), float(w.get("x0", 0.0))))
    blocks: list[dict] = []
    cur: list[dict] = []

    def flush() -> None:
        if not cur:
            return
        spans = []
        for w in cur:
            font = w.get("fontname", "") or ""
            flags = _BOLD_FLAG if "bold" in font.lower() else 0
            spans.append({
                "bbox": (float(w["x0"]), float(w["top"]), float(w["x1"]), float(w["bottom"])),
                "text": w.get("text") or "",
                "size": float(w.get("size", 0) or 0),
                "font": font,
                "flags": flags,
            })
        bx0 = min(s["bbox"][0] for s in spans); by0 = min(s["bbox"][1] for s in spans)
        bx1 = max(s["bbox"][2] for s in spans); by1 = max(s["bbox"][3] for s in spans)
        blocks.append({"type": 0, "bbox": (bx0, by0, bx1, by1), "lines": [{"spans": spans}]})
        cur.clear()

    for w in words:
        if cur and abs(float(w.get("top", 0.0)) - float(cur[0].get("top", 0.0))) <= 3.0:
            cur.append(w)
        else:
            flush()
            cur.append(w)
    flush()
    return blocks


def _parse_pdf(path: str) -> str:
    import pdfplumber

    doc = pdfplumber.open(path)
    out: list[str] = []
    for page in doc.pages:
        blocks = _page_to_blocks(page)
        words_xy = [(float(w["x0"]), float(w["top"]), float(w["x1"]), float(w["bottom"]), w.get("text") or "")
                    for w in page.extract_words()]
        tables = _ruled_tables(page, words_xy) + _borderless_regions(blocks)
        tregions = [t["bbox"] for t in tables]
        sizes = [s["size"] for b in blocks for ln in b.get("lines", []) for s in ln.get("spans", [])]
        body = statistics.median(sizes) if sizes else 10.0
        ordered = _reading_order(blocks, float(page.width))
        placed: set[int] = set()
        page_md: list[str] = []
        for b in ordered:
            ti = next((i for i, r in enumerate(tregions) if _inside(b["bbox"], r)), None)
            if ti is not None:
                if ti not in placed:
                    page_md.append(tables[ti]["md"])
                    placed.add(ti)
                continue
            for ln in b.get("lines", []):
                spans = ln.get("spans", [])
                txt = " ".join(s["text"] for s in spans).strip()
                if not txt:
                    continue
                if _line_is_formula(spans):
                    page_md.append(f"$$ {txt} $$")   # Formula block
                    continue
                mx = max((s["size"] for s in spans), default=body)
                bold = any(s.get("flags", 0) & 16 for s in spans)
                if mx >= body * 1.5:
                    page_md.append(f"# {txt}")
                elif mx >= body * 1.25 or (bold and len(txt) < 40):
                    page_md.append(f"## {txt}")
                else:
                    page_md.append(txt)
        for i, t in enumerate(tables):
            if i not in placed:
                page_md.append(t["md"])
        out.append("\n\n".join(page_md))
    doc.close()
    return "\n\n".join(out).strip()


def _parse_docx(path: str) -> str:
    import docx

    d = docx.Document(path)
    out: list[str] = []
    for p in d.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        style = (p.style.name or "").lower() if p.style else ""
        if "heading 1" in style or style == "title":
            out.append(f"# {t}")
        elif "heading" in style:
            out.append(f"## {t}")
        else:
            out.append(t)
    for tbl in d.tables:
        rows = [[c.text.strip() for c in r.cells] for r in tbl.rows]
        if rows:
            out.append(_rows_to_md(rows))
    return "\n\n".join(out).strip()


def _parse_xlsx(path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        rows = [[("" if c is None else str(c)).strip() for c in row]
                for row in ws.iter_rows(values_only=True)]
        rows = [r for r in rows if any(r)]
        if rows:
            out.append(f"## {ws.title}")
            out.append(_rows_to_md(rows))
    wb.close()
    return "\n\n".join(out).strip()


def _parse_pptx(path: str) -> str:
    import pptx

    prs = pptx.Presentation(path)
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        out.append(t)
    return "\n\n".join(out).strip()


def parse(path: str, mime: str) -> str:
    """Parse a document into structured Markdown. Raises ValueError for unsupported types."""
    kind = SUPPORTED.get(mime)
    if kind == "pdf":
        return _parse_pdf(path)
    if kind == "docx":
        return _parse_docx(path)
    if kind == "xlsx":
        return _parse_xlsx(path)
    if kind == "pptx":
        return _parse_pptx(path)
    raise ValueError(f"unsupported mime: {mime}")
